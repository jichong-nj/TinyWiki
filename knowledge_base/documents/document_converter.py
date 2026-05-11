import os
import json
import tempfile
import subprocess
import shutil
from pathlib import Path
from django.utils import timezone
from django.conf import settings
from .storage_models import FileStorage, FileAttachment, DocumentConversion
from .storage_service import StorageService


class DocumentConverter:
    """文档转换器"""
    
    SUPPORTED_TYPES = ['docx', 'doc', 'md', 'txt', 'pptx', 'ppt', 'pdf']
    
    @classmethod
    def convert(cls, file_storage_id):
        """
        同步转换文档
        
        Args:
            file_storage_id: FileStorage对象的ID
            
        Returns:
            dict: 转换结果
        """
        try:
            file_storage = StorageService.get_file_by_id(file_storage_id)
            if not file_storage:
                return {'success': False, 'error': '文件不存在'}
            
            # 创建转换记录
            conversion = DocumentConversion.objects.create(
                original_file=file_storage,
                status='converting',
                started_at=timezone.now()
            )
            
            # 根据文件类型选择转换器
            file_type = file_storage.file_type.lower()
            
            if file_type in ['docx']:
                result = cls._convert_docx(file_storage, conversion)
            elif file_type in ['pptx', 'ppt']:
                result = cls._convert_pptx(file_storage, conversion)
            elif file_type in ['md', 'txt']:
                result = cls._convert_text(file_storage, conversion)
            elif file_type == 'pdf':
                result = cls._convert_pdf(file_storage, conversion)
            else:
                result = {'success': False, 'error': f'不支持的文件类型: {file_type}'}
            
            # 更新转换状态
            if result['success']:
                conversion.status = 'success'
                conversion.converted_file = result.get('converted_file')
                conversion.conversion_info = result.get('info', {})
            else:
                conversion.status = 'failed'
                conversion.error_message = result.get('error', '未知错误')
            
            conversion.completed_at = timezone.now()
            conversion.save()
            
            return result
            
        except Exception as e:
            import traceback
            return {'success': False, 'error': f'转换失败: {str(e)}\n{traceback.format_exc()}'}
    
    @classmethod
    def convert_async(cls, file_storage_id):
        """
        异步转换文档（同步实现的兼容方法）
        
        Args:
            file_storage_id: FileStorage对象的ID
            
        Returns:
            dict: 转换结果
        """
        # 目前直接调用同步方法
        return cls.convert(file_storage_id)
    
    @classmethod
    def _convert_docx(cls, file_storage, conversion):
        """
        转换Word文档
        
        Args:
            file_storage: FileStorage对象
            conversion: DocumentConversion对象
            
        Returns:
            dict: 转换结果
        """
        try:
            # 尝试使用python-docx
            import docx
        except ImportError:
            return cls._convert_docx_simple(file_storage, conversion)
        
        try:
            document = docx.Document(file_storage.full_path)
            
            # 存储提取的图片
            attachments = []
            image_map = {}
            
            # 提取图片
            for rel in document.part.rels.values():
                if 'image' in rel.target_ref:
                    image_data = rel.target_part.blob
                    
                    # 保存图片为附件
                    image_name = f"image_{len(attachments)+1}.{rel.target_ref.split('.')[-1]}"
                    
                    # 创建临时文件
                    with tempfile.NamedTemporaryFile(delete=False, suffix=f".{rel.target_ref.split('.')[-1]}") as f:
                        f.write(image_data)
                        temp_path = f.name
                    
                    try:
                        # 保存到附件库
                        image_file = StorageService.save_file_from_path(
                            source_path=temp_path,
                            storage_type='attachment',
                            document=file_storage.document,
                            original_name=image_name
                        )
                        
                        # 创建附件关联
                        if file_storage.document:
                            FileAttachment.objects.create(
                                document=file_storage.document,
                                file_storage=image_file,
                                original_name=image_name,
                                description=f"从{file_storage.file_name}中提取的图片",
                                position=len(attachments)
                            )
                        
                        attachments.append(image_file)
                        image_map[rel.rId] = image_file.file_url
                    finally:
                        # 删除临时文件
                        os.unlink(temp_path)
            
            # 转换内容为Markdown
            markdown_content = []
            
            # 处理段落（包括段落内的图片）
            for paragraph in document.paragraphs:
                # 检查段落中是否有图片并插入
                for run in paragraph.runs:
                    for inline in run.element.findall('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                        for blip in inline.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip'):
                            rId = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                            if rId and rId in image_map:
                                # 找到对应的图片文件
                                img_file = None
                                for idx, af in enumerate(attachments):
                                    if af.file_url == image_map[rId]:
                                        img_file = af
                                        break
                                if img_file:
                                    markdown_content.append(f"![{img_file.file_name}]({image_map[rId]})\n")
                
                # 不管有没有图片，都处理文本
                if paragraph.style.name.startswith('Heading 1'):
                    markdown_content.append(f"# {paragraph.text}\n")
                elif paragraph.style.name.startswith('Heading 2'):
                    markdown_content.append(f"## {paragraph.text}\n")
                elif paragraph.style.name.startswith('Heading 3'):
                    markdown_content.append(f"### {paragraph.text}\n")
                elif paragraph.text.strip():
                    markdown_content.append(f"{paragraph.text}\n")
                else:
                    markdown_content.append("\n")
            
            # 处理表格
            for table in document.tables:
                markdown_content.append("\n")
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    markdown_content.append(f"| {' | '.join(cells)} |\n")
                    if i == 0:
                        markdown_content.append(f"| {' | '.join(['---'] * len(cells))} |\n")
            
            # 如果有图片但没有在段落中找到，就全部添加到最后
            if len(attachments) > 0 and not any('![' in item for item in markdown_content):
                markdown_content.append("\n---\n")
                markdown_content.append("## 提取的图片\n\n")
                for img_file in attachments:
                    markdown_content.append(f"![{img_file.file_name}]({img_file.file_url})\n\n")
            
            markdown_text = ''.join(markdown_content)
            
            # 保存转换后的Markdown
            converted_name = Path(file_storage.file_name).stem + '.md'
            converted_file = StorageService.save_file(
                file_obj=None,
                storage_type='converted',
                document=file_storage.document,
                original_name=converted_name,
                content=markdown_text
            )
            
            # 更新文档内容（如果有关联文档）
            if file_storage.document:
                file_storage.document.content = markdown_text
                file_storage.document.save()
            
            # 保存转换信息
            conversion_info = {
                'attachments_count': len(attachments),
                'image_map': {k: v for k, v in image_map.items()},
                'paragraphs_count': len(document.paragraphs),
                'tables_count': len(document.tables)
            }
            
            return {
                'success': True,
                'converted_file': converted_file,
                'info': conversion_info,
                'markdown': markdown_text
            }
            
        except Exception as e:
            import traceback
            return {'success': False, 'error': f'DOCX转换失败: {str(e)}\n{traceback.format_exc()}'}
    
    @classmethod
    def _convert_docx_simple(cls, file_storage, conversion):
        """
        简化的Word文档转换（当没有python-docx时使用）
        
        Args:
            file_storage: FileStorage对象
            conversion: DocumentConversion对象
            
        Returns:
            dict: 转换结果
        """
        try:
            # 读取原始文件内容
            content = f"# {Path(file_storage.file_name).stem}\n\n"
            content += f"此文件为Word文档，请安装python-docx库以获得更好的转换效果。\n\n"
            content += f"原始文件名: {file_storage.file_name}\n"
            content += f"文件大小: {file_storage.file_size} bytes\n"
            
            # 保存转换后的Markdown
            converted_name = Path(file_storage.file_name).stem + '.md'
            converted_file = StorageService.save_file(
                file_obj=None,
                storage_type='converted',
                document=file_storage.document,
                original_name=converted_name,
                content=content
            )
            
            # 更新文档内容（如果有关联文档）
            if file_storage.document:
                file_storage.document.content = content
                file_storage.document.save()
            
            conversion_info = {
                'method': 'simple',
                'warning': '未安装python-docx库，使用简化转换'
            }
            
            return {
                'success': True,
                'converted_file': converted_file,
                'info': conversion_info,
                'markdown': content
            }
            
        except Exception as e:
            return {'success': False, 'error': f'简化转换失败: {str(e)}'}
    
    @classmethod
    def _convert_text(cls, file_storage, conversion):
        """
        转换文本文件（Markdown或TXT）
        
        Args:
            file_storage: FileStorage对象
            conversion: DocumentConversion对象
            
        Returns:
            dict: 转换结果
        """
        try:
            # 直接读取文件内容
            with open(file_storage.full_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # 如果是txt文件，添加标题
            if file_storage.file_type.lower() == 'txt':
                title = Path(file_storage.file_name).stem
                content = f"# {title}\n\n{content}"
            
            # 保存转换后的Markdown（其实就是原样保存）
            converted_name = Path(file_storage.file_name).stem + '.md'
            converted_file = StorageService.save_file(
                file_obj=None,
                storage_type='converted',
                document=file_storage.document,
                original_name=converted_name,
                content=content
            )
            
            # 更新文档内容（如果有关联文档）
            if file_storage.document:
                file_storage.document.content = content
                file_storage.document.save()
            
            conversion_info = {
                'method': 'direct',
                'content_length': len(content)
            }
            
            return {
                'success': True,
                'converted_file': converted_file,
                'info': conversion_info,
                'markdown': content
            }
            
        except Exception as e:
            return {'success': False, 'error': f'文本转换失败: {str(e)}'}
    
    @classmethod
    def _convert_pptx(cls, file_storage, conversion):
        """
        转换PowerPoint文档
        
        Args:
            file_storage: FileStorage对象
            conversion: DocumentConversion对象
            
        Returns:
            dict: 转换结果
        """
        try:
            # 尝试导入必需的库
            from pptx import Presentation
            from pdf2image import convert_from_path
            from PIL import Image
            
        except ImportError as e:
            return {
                'success': False, 
                'error': f'缺少依赖库: {str(e)}\n请安装: pip install python-pptx pdf2image Pillow'
            }
        
        try:
            ppt_path = file_storage.full_path
            
            # 1. 读取PPT文件，提取备注
            prs = Presentation(ppt_path)
            slide_info_list = []
            
            for idx, slide in enumerate(prs.slides):
                # 提取备注
                note_text = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    for shape in notes_slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            note_text += shape.text + "\n"
                
                # 提取每页文字
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        try:
                            slide_text += shape.text + "\n"
                        except:
                            pass
                
                slide_info_list.append({
                    'index': idx + 1,
                    'text': slide_text.strip(),
                    'notes': note_text.strip()
                })
            
            # 2. 尝试 PPT → PDF 转换
            pdf_path = None
            temp_dir = tempfile.mkdtemp()
            try:
                # 使用 LibreOffice 转换
                output_dir = temp_dir
                result = subprocess.run([
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', output_dir,
                    ppt_path
                ], capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    # 找到生成的 PDF 文件
                    pdf_filename = Path(ppt_path).stem + '.pdf'
                    pdf_path = Path(output_dir) / pdf_filename
                    
                    if not pdf_path.exists():
                        # 尝试找其他可能的文件名
                        pdf_files = list(Path(output_dir).glob("*.pdf"))
                        if pdf_files:
                            pdf_path = pdf_files[0]
            except Exception as e:
                pdf_path = None
            
            # 3. PDF → 图片 转换
            slide_images = []
            attachments = []
            
            if pdf_path and pdf_path.exists():
                try:
                    images = convert_from_path(
                        str(pdf_path),
                        dpi=300,
                        fmt='png'
                    )
                    
                    for idx, img in enumerate(images):
                        # 保存图片到临时文件
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                            img.save(tmp, format='PNG')
                            tmp_path = tmp.name
                        
                        # 保存到 storage/attachments
                        img_filename = f"slide_{idx+1}.png"
                        img_file = StorageService.save_file_from_path(
                            source_path=tmp_path,
                            storage_type='attachment',
                            document=file_storage.document,
                            original_name=img_filename
                        )
                        
                        slide_images.append({
                            'index': idx + 1,
                            'file': img_file
                        })
                        
                        attachments.append(img_file)
                        
                        # 清理临时文件
                        os.unlink(tmp_path)
                except Exception as e:
                    slide_images = []
            
            # 4. 生成 Markdown 内容
            markdown_content = []
            
            # 标题
            markdown_content.append(f"# {Path(file_storage.file_name).stem}\n\n")
            
            for idx, slide_info in enumerate(slide_info_list):
                slide_num = slide_info['index']
                
                # 添加幻灯片标题
                markdown_content.append(f"## 第 {slide_num} 页\n\n")
                
                # 添加图片
                img_found = False
                for slide_img in slide_images:
                    if slide_img['index'] == slide_num:
                        markdown_content.append(f"![幻灯片 {slide_num}]({slide_img['file'].file_url})\n\n")
                        img_found = True
                        break
                
                if not img_found:
                    markdown_content.append(f"![幻灯片 {slide_num}](未生成图片)\n\n")
                
                # 添加文字内容
                if slide_info['text']:
                    markdown_content.append(f"**内容：**\n\n{slide_info['text']}\n\n")
                
                # 添加备注
                if slide_info['notes']:
                    markdown_content.append(f"**备注：**\n\n{slide_info['notes']}\n\n")
                
                markdown_content.append("---\n\n")
            
            # 保存转换后的Markdown
            converted_name = Path(file_storage.file_name).stem + '.md'
            markdown_text = ''.join(markdown_content)
            
            converted_file = StorageService.save_file(
                file_obj=None,
                storage_type='converted',
                document=file_storage.document,
                original_name=converted_name,
                content=markdown_text
            )
            
            # 更新文档内容
            if file_storage.document:
                file_storage.document.content = markdown_text
                file_storage.document.save()
            
            conversion_info = {
                'slides_count': len(slide_info_list),
                'images_count': len(slide_images),
                'method': 'libreoffice-pdf' if pdf_path else 'text-only'
            }
            
            # 清理临时目录
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
            
            return {
                'success': True,
                'converted_file': converted_file,
                'info': conversion_info,
                'markdown': markdown_text
            }
            
        except Exception as e:
            import traceback
            return {'success': False, 'error': f'PPT转换失败: {str(e)}\n{traceback.format_exc()}'}
    
    @classmethod
    def _extract_pdf_text_with_pypdf2(cls, pdf_path):
        """
        使用 PyPDF2 提取 PDF 文本
        
        Args:
            pdf_path: PDF 文件路径
            
        Returns:
            tuple: (是否成功, 文本内容, 提取方法)
        """
        try:
            from PyPDF2 import PdfReader
            
            reader = PdfReader(pdf_path)
            text_parts = []
            
            for page_num, page in enumerate(reader.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(f"\n--- 第 {page_num + 1} 页 ---\n")
                    text_parts.append(page_text)
            
            full_text = ''.join(text_parts).strip()
            if full_text and len(full_text) > 100:
                return True, full_text, 'pypdf2-text'
            return False, '', 'no-text'
            
        except Exception as e:
            return False, '', f'error: {str(e)}'
    
    @classmethod
    def _convert_pdf_with_llm(cls, pdf_path, file_storage, conversion):
        """
        将 PDF 转为图片，并用 LLM 提取文字
        
        Args:
            pdf_path: PDF 文件路径
            file_storage: FileStorage 对象
            conversion: DocumentConversion 对象
            
        Returns:
            dict: 转换结果
        """
        try:
            from pdf2image import convert_from_path
            from PIL import Image
            
            # 1. 将 PDF 转成图片
            temp_dir = tempfile.mkdtemp()
            try:
                images = convert_from_path(
                    str(pdf_path),
                    dpi=300,
                    fmt='png'
                )
                
                # 2. 处理每一页图片
                page_info_list = []
                attachments = []
                markdown_content = []
                
                # 标题
                markdown_content.append(f"# {Path(file_storage.file_name).stem}\n\n")
                
                for idx, img in enumerate(images):
                    page_num = idx + 1
                    
                    # 保存图片到临时文件
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                        img.save(tmp, format='PNG')
                        tmp_path = tmp.name
                    
                    try:
                        # 保存到 storage/attachments
                        img_filename = f"page_{page_num}.png"
                        img_file = StorageService.save_file_from_path(
                            source_path=tmp_path,
                            storage_type='attachment',
                            document=file_storage.document,
                            original_name=img_filename
                        )
                        attachments.append(img_file)
                        
                        # 简单的 OCR/LLM 占位符
                        # 可以后续添加真正的 OCR 或 LLM 提取功能
                        extracted_text = f"（第 {page_num} 页 - 图片已保存\n"
                        
                        # 保存页面信息
                        page_info_list.append({
                            'index': page_num,
                            'image_file': img_file,
                            'text': extracted_text
                        })
                        
                        # 添加到 Markdown
                        markdown_content.append(f"## 第 {page_num} 页\n\n")
                        markdown_content.append(f"![第 {page_num} 页]({img_file.file_url})\n\n")
                        markdown_content.append(f"{extracted_text}\n\n")
                        
                    finally:
                        os.unlink(tmp_path)
                
                # 生成 Markdown
                markdown_text = ''.join(markdown_content)
                
                # 保存转换后的 Markdown
                converted_name = Path(file_storage.file_name).stem + '.md'
                converted_file = StorageService.save_file(
                    file_obj=None,
                    storage_type='converted',
                    document=file_storage.document,
                    original_name=converted_name,
                    content=markdown_text
                )
                
                # 更新文档内容
                if file_storage.document:
                    file_storage.document.content = markdown_text
                    file_storage.document.save()
                
                conversion_info = {
                    'method': 'image-only',
                    'pages_count': len(page_info_list),
                    'images_count': len(attachments),
                    'note': 'PDF 扫描版：图片已保存'
                }
                
                # 清理临时目录
                try:
                    shutil.rmtree(temp_dir)
                except:
                    pass
                
                return {
                    'success': True,
                    'converted_file': converted_file,
                    'info': conversion_info,
                    'markdown': markdown_text
                }
                
            except Exception as e:
                try:
                    shutil.rmtree(temp_dir)
                except:
                    pass
                raise
                
        except Exception as e:
            import traceback
            return {'success': False, 'error': f'PDF 图片转换失败: {str(e)}\n{traceback.format_exc()}'}
    
    @classmethod
    def _convert_pdf(cls, file_storage, conversion):
        """
        转换 PDF 文档
        
        Args:
            file_storage: FileStorage 对象
            conversion: DocumentConversion 对象
            
        Returns:
            dict: 转换结果
        """
        try:
            pdf_path = file_storage.full_path
            
            # 1. 先尝试用 PyPDF2 提取文本
            success, extracted_text, method = cls._extract_pdf_text_with_pypdf2(pdf_path)
            
            # 2. 如果有文本，按可解析的 PDF 处理（类似 Word）
            if success:
                try:
                    from PyPDF2 import PdfReader
                    
                    reader = PdfReader(pdf_path)
                    
                    # 存储提取的图片
                    attachments = []
                    image_map = {}
                    
                    # 尝试提取图片
                    for page_num, page in enumerate(reader.pages):
                        if '/Resources' in page and '/XObject' in page['/Resources']:
                            xObject = page['/Resources']['/XObject'].get_object()
                            for obj in xObject:
                                if xObject[obj]['/Subtype'] == '/Image':
                                    try:
                                        data = xObject[obj]._data
                                        img_ext = 'png'
                                        
                                        with tempfile.NamedTemporaryFile(delete=False, suffix=f'.{img_ext}') as f:
                                            f.write(data)
                                            temp_path = f.name
                                        
                                        try:
                                            image_name = f"page{page_num+1}_image{len(attachments)+1}.{img_ext}"
                                            image_file = StorageService.save_file_from_path(
                                                source_path=temp_path,
                                                storage_type='attachment',
                                                document=file_storage.document,
                                                original_name=image_name
                                            )
                                            
                                            if file_storage.document:
                                                FileAttachment.objects.create(
                                                    document=file_storage.document,
                                                    file_storage=image_file,
                                                    original_name=image_name,
                                                    description=f"从 PDF 第 {page_num+1} 页提取的图片",
                                                    position=len(attachments)
                                                )
                                            
                                            attachments.append(image_file)
                                        finally:
                                            os.unlink(temp_path)
                                    except Exception as img_e:
                                        continue
                    
                    # 生成 Markdown
                    markdown_content = []
                    markdown_content.append(f"# {Path(file_storage.file_name).stem}\n\n")
                    markdown_content.append(extracted_text)
                    
                    # 如果有提取的图片，添加到最后
                    if len(attachments) > 0:
                        markdown_content.append("\n---\n")
                        markdown_content.append("## 提取的图片\n\n")
                        for img_file in attachments:
                            markdown_content.append(f"![{img_file.file_name}]({img_file.file_url})\n\n")
                    
                    markdown_text = ''.join(markdown_content)
                    
                    # 保存转换后的 Markdown
                    converted_name = Path(file_storage.file_name).stem + '.md'
                    converted_file = StorageService.save_file(
                        file_obj=None,
                        storage_type='converted',
                        document=file_storage.document,
                        original_name=converted_name,
                        content=markdown_text
                    )
                    
                    # 更新文档内容
                    if file_storage.document:
                        file_storage.document.content = markdown_text
                        file_storage.document.save()
                    
                    conversion_info = {
                        'method': f'text-based-{method}',
                        'attachments_count': len(attachments),
                        'pages_count': len(reader.pages)
                    }
                    
                    return {
                        'success': True,
                        'converted_file': converted_file,
                        'info': conversion_info,
                        'markdown': markdown_text
                    }
                    
                except Exception as e:
                    # 如果 PyPDF2 处理失败，回退到 LLM 图片模式
                    import traceback
                    print(f"PyPDF2 解析失败，回退到图片模式: {e}\n{traceback.format_exc()}")
            
            # 3. 没有文本或者 PyPDF2 失败，用图片 + LLM 模式
            return cls._convert_pdf_with_llm(pdf_path, file_storage, conversion)
            
        except Exception as e:
            import traceback
            return {'success': False, 'error': f'PDF 转换失败: {str(e)}\n{traceback.format_exc()}'}
